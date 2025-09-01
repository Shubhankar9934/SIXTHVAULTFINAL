"""
ENHANCED MULTIMODAL TEXT EXTRACTOR v2.0

Lightweight text extractor for PDF, DOCX, Excel, PPT, Images, and more.
Now supports multimodal content extraction with backward compatibility.

Features:
- Original lightweight extraction (backward compatible)
- Enhanced multimodal extraction with metadata
- OCR support for images
- Table structure preservation
- PowerPoint slide extraction
- Fallback strategies
"""
from pathlib import Path
import pdfplumber, docx2txt, pandas as pd
import time
from typing import Dict, List, Any, Optional
from dataclasses import dataclass

# Enhanced extraction libraries
try:
    import fitz  # PyMuPDF for advanced PDF processing
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("OCR not available - image text extraction disabled")

try:
    from docx import Document as DocxDocument
    DOCX_ADVANCED_AVAILABLE = True
except ImportError:
    DOCX_ADVANCED_AVAILABLE = False

@dataclass
class ExtractedContent:
    """Container for extracted multimodal content"""
    text: str
    images: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    content_type: str
    file_path: str
    extraction_time: float
    
    def __post_init__(self):
        if self.images is None:
            self.images = []
        if self.tables is None:
            self.tables = []
        if self.metadata is None:
            self.metadata = {}
    
    def get_combined_text(self) -> str:
        """Get all text content combined for backward compatibility"""
        combined = [self.text]
        
        # Add table descriptions
        for table in self.tables:
            if table.get('description'):
                combined.append(f"[TABLE] {table['description']}")
        
        # Add image descriptions and OCR text
        for image in self.images:
            if image.get('description'):
                combined.append(f"[IMAGE] {image['description']}")
            if image.get('ocr_text'):
                combined.append(f"[OCR] {image['ocr_text']}")
        
        return "\n\n".join(filter(None, combined))

def extract_text(file_path: str) -> str:
    """
    BACKWARD COMPATIBLE: Original lightweight text extraction
    This function maintains full backward compatibility with existing code
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        with pdfplumber.open(path) as pdf:
            return "\n".join((p.extract_text() or "") for p in pdf.pages)

    if ext in {".docx", ".doc"}:
        return docx2txt.process(path)

    if ext in {".xls", ".xlsx"}:
        dfs = pd.read_excel(path, sheet_name=None)
        return "\n".join(df.to_csv(index=False) for df in dfs.values())

    # fall back to plain text
    return path.read_text(encoding="utf-8", errors="ignore")

def extract_content_multimodal(file_path: str) -> ExtractedContent:
    """
    ENHANCED: Multimodal content extraction with metadata
    Use this for new implementations that need multimodal data
    """
    start_time = time.time()
    path = Path(file_path)
    ext = path.suffix.lower()
    
    # Initialize content container
    text_content = ""
    images = []
    tables = []
    content_type = "unknown"
    
    try:
        if ext == ".pdf":
            text_content, images, tables = _extract_pdf_multimodal(file_path)
            content_type = "pdf_document"
        
        elif ext in {".docx", ".doc"}:
            text_content, images, tables = _extract_word_multimodal(file_path)
            content_type = "word_document"
        
        elif ext in {".xls", ".xlsx"}:
            text_content, tables = _extract_excel_multimodal(file_path)
            content_type = "excel_spreadsheet"
        
        elif ext in {".pptx", ".ppt"} and PPTX_AVAILABLE:
            text_content, images, tables = _extract_ppt_multimodal(file_path)
            content_type = "powerpoint_presentation"
        
        elif ext in {".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".gif"} and OCR_AVAILABLE:
            text_content, images = _extract_image_multimodal(file_path)
            content_type = "image_file"
        
        else:
            # Fallback to plain text
            text_content = path.read_text(encoding="utf-8", errors="ignore")
            content_type = "text_file"
    
    except Exception as e:
        print(f"Multimodal extraction error for {file_path}: {e}")
        # Fallback to basic extraction
        text_content = extract_text(file_path)
        content_type = "fallback"
    
    # Create metadata
    stat = path.stat()
    metadata = {
        "filename": path.name,
        "file_size": stat.st_size,
        "file_extension": ext,
        "created_time": stat.st_ctime,
        "modified_time": stat.st_mtime,
        "file_type": _detect_file_type(ext)
    }
    
    extraction_time = time.time() - start_time
    
    return ExtractedContent(
        text=text_content,
        images=images,
        tables=tables,
        metadata=metadata,
        content_type=content_type,
        file_path=file_path,
        extraction_time=extraction_time
    )

def _extract_pdf_multimodal(file_path: str) -> tuple:
    """Extract PDF with images and tables"""
    text_content = []
    images = []
    tables = []
    
    try:
        # Primary extraction with pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # Extract text
                page_text = page.extract_text()
                if page_text:
                    text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                
                # Extract tables
                page_tables = page.extract_tables()
                for table_idx, table in enumerate(page_tables):
                    if table and len(table) > 1:
                        try:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            tables.append({
                                'page': page_num + 1,
                                'table_index': table_idx,
                                'description': f"Table with {len(df)} rows and {len(df.columns)} columns from page {page_num + 1}",
                                'csv_representation': df.to_csv(index=False)
                            })
                        except Exception as table_error:
                            print(f"Table extraction error: {table_error}")
        
        # Advanced image extraction with PyMuPDF if available
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    image_list = page.get_images()
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            if pix.n - pix.alpha < 4:  # GRAY or RGB
                                img_data = pix.tobytes("png")
                                
                                # Perform OCR if available
                                ocr_text = ""
                                if OCR_AVAILABLE:
                                    try:
                                        from io import BytesIO
                                        image = Image.open(BytesIO(img_data))
                                        ocr_text = pytesseract.image_to_string(image).strip()
                                    except Exception:
                                        pass
                                
                                images.append({
                                    'page': page_num + 1,
                                    'image_index': img_index,
                                    'format': 'png',
                                    'size': len(img_data),
                                    'ocr_text': ocr_text,
                                    'description': f"Image from page {page_num + 1}"
                                })
                            
                            pix = None
                        except Exception as img_error:
                            print(f"Image extraction error: {img_error}")
                doc.close()
            except Exception as pymupdf_error:
                print(f"PyMuPDF extraction error: {pymupdf_error}")
    
    except Exception as e:
        print(f"PDF multimodal extraction error: {e}")
        # Fallback to basic extraction
        text_content = [extract_text(file_path)]
    
    return "\n\n".join(filter(None, text_content)), images, tables

def _extract_word_multimodal(file_path: str) -> tuple:
    """Extract Word document with tables"""
    text_content = ""
    images = []
    tables = []
    
    try:
        if DOCX_ADVANCED_AVAILABLE and Path(file_path).suffix.lower() == '.docx':
            # Advanced extraction with python-docx
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text_content = "\n\n".join(paragraphs)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                try:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data and len(table_data) > 1:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        tables.append({
                            'table_index': table_idx,
                            'description': f"Table with {len(df)} rows and {len(df.columns)} columns",
                            'csv_representation': df.to_csv(index=False)
                        })
                except Exception as table_error:
                    print(f"Word table extraction error: {table_error}")
            
            # Note presence of images (extraction requires additional libraries)
            if any(shape.shape_type == 13 for shape in doc.inline_shapes):  # Picture type
                images.append({'description': 'Images present in document'})
        else:
            # Fallback to docx2txt
            text_content = docx2txt.process(file_path)
    
    except Exception as e:
        print(f"Word multimodal extraction error: {e}")
        text_content = extract_text(file_path)
    
    return text_content, images, tables

def _extract_excel_multimodal(file_path: str) -> tuple:
    """Extract Excel with enhanced table analysis"""
    text_content = []
    tables = []
    
    try:
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Generate text description
            sheet_description = f"[Sheet: {sheet_name}]\n"
            sheet_description += f"Dimensions: {df.shape[0]} rows × {df.shape[1]} columns\n"
            sheet_description += f"Columns: {', '.join(df.columns.astype(str))}\n"
            
            # Add sample data
            if not df.empty:
                sheet_description += "\nSample data:\n"
                sheet_description += df.head(3).to_string(index=False)
            
            text_content.append(sheet_description)
            
            # Store table data
            tables.append({
                'sheet_name': sheet_name,
                'description': f"Excel sheet '{sheet_name}' with {df.shape[0]} rows and {df.shape[1]} columns",
                'csv_representation': df.to_csv(index=False)
            })
    
    except Exception as e:
        print(f"Excel multimodal extraction error: {e}")
        text_content = [extract_text(file_path)]
    
    return "\n\n".join(text_content), tables

def _extract_ppt_multimodal(file_path: str) -> tuple:
    """Extract PowerPoint with slides and tables"""
    text_content = []
    images = []
    tables = []
    
    try:
        prs = Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = f"[Slide {slide_idx + 1}]\n"
            slide_content = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # Check for tables
                if shape.has_table:
                    try:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data and len(table_data) > 1:
                            df = pd.DataFrame(table_data[1:], columns=table_data[0])
                            tables.append({
                                'slide': slide_idx + 1,
                                'description': f"Table from slide {slide_idx + 1}",
                                'csv_representation': df.to_csv(index=False)
                            })
                    except Exception as table_error:
                        print(f"PPT table extraction error: {table_error}")
            
            if slide_content:
                slide_text += "\n".join(slide_content)
                text_content.append(slide_text)
            
            # Count images
            image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # Picture type
            if image_count > 0:
                images.append({
                    'slide': slide_idx + 1,
                    'count': image_count,
                    'description': f"{image_count} images on slide {slide_idx + 1}"
                })
    
    except Exception as e:
        print(f"PowerPoint multimodal extraction error: {e}")
        text_content = ["Error extracting PowerPoint content"]
    
    return "\n\n".join(text_content), images, tables

def _extract_image_multimodal(file_path: str) -> tuple:
    """Extract image with OCR"""
    text_content = ""
    images = []
    
    try:
        # Perform OCR on the image
        image = Image.open(file_path)
        ocr_text = pytesseract.image_to_string(image)
        
        text_content = f"[OCR Text from {Path(file_path).name}]\n{ocr_text.strip()}"
        
        # Store image metadata
        images.append({
            'filename': Path(file_path).name,
            'format': image.format,
            'size': image.size,
            'mode': image.mode,
            'ocr_text': ocr_text.strip(),
            'description': f"Image file with OCR text extraction"
        })
    
    except Exception as e:
        print(f"Image extraction error: {e}")
        text_content = f"Error extracting text from image: {Path(file_path).name}"
    
    return text_content, images

def _detect_file_type(ext: str) -> str:
    """Detect file type from extension"""
    type_mapping = {
        '.pdf': 'document',
        '.docx': 'document', '.doc': 'document',
        '.xlsx': 'spreadsheet', '.xls': 'spreadsheet',
        '.pptx': 'presentation', '.ppt': 'presentation',
        '.jpg': 'image', '.jpeg': 'image', '.png': 'image', 
        '.tiff': 'image', '.bmp': 'image', '.gif': 'image',
        '.txt': 'text', '.rtf': 'text', '.md': 'text'
    }
    return type_mapping.get(ext, 'unknown')

def get_supported_formats() -> List[str]:
    """Get list of supported file formats"""
    formats = ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.txt', '.rtf', '.md']
    
    if PPTX_AVAILABLE:
        formats.extend(['.pptx', '.ppt'])
    
    if OCR_AVAILABLE:
        formats.extend(['.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.gif'])
    
    return sorted(formats)

def get_extraction_capabilities() -> Dict[str, Any]:
    """Get information about extraction capabilities"""
    return {
        'basic_extraction': True,
        'multimodal_extraction': True,
        'pdf_advanced': PYMUPDF_AVAILABLE,
        'powerpoint_support': PPTX_AVAILABLE,
        'ocr_support': OCR_AVAILABLE,
        'docx_advanced': DOCX_ADVANCED_AVAILABLE,
        'supported_formats': get_supported_formats()
    }
