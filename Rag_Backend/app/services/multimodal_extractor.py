"""
ADVANCED MULTIMODAL FILE FORMAT HANDLER v1.0
==============================================

Common file format handling and extraction methodology for:
- PDF Documents (text + images + tables)
- Word Documents (DOCX/DOC)
- Excel Spreadsheets (XLS/XLSX)
- PowerPoint Presentations (PPT/PPTX)
- Images (JPEG, PNG, TIFF, BMP)
- Text Files (TXT, RTF, MD)

Features:
- Unified extraction interface
- Multimodal content preservation
- Metadata extraction
- OCR for image-based content
- Table structure preservation
- Slide content extraction
- Error handling and fallbacks
"""

from __future__ import annotations
import os
import time
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple, Union
from dataclasses import dataclass
from abc import ABC, abstractmethod
import logging

# Core extraction libraries
import pdfplumber
import docx2txt
import pandas as pd
from PIL import Image
import pytesseract

# Advanced extraction libraries
try:
    import fitz  # PyMuPDF for advanced PDF processing
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("PyMuPDF not available - using pdfplumber only")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx not available - PPT extraction disabled")

try:
    from docx import Document as DocxDocument
    DOCX_ADVANCED_AVAILABLE = True
except ImportError:
    DOCX_ADVANCED_AVAILABLE = False
    print("python-docx not available - using docx2txt only")

# ═══════════════════════════════════════════════════════════════════════════
# Data Structures for Multimodal Content
# ═══════════════════════════════════════════════════════════════════════════

@dataclass
class ExtractedContent:
    """Container for extracted multimodal content"""
    text: str
    images: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    content_type: str
    file_path: str
    extraction_time: float
    
    def get_combined_text(self) -> str:
        """Get all text content combined"""
        combined = [self.text]
        
        # Add table descriptions
        for table in self.tables:
            if table.get('description'):
                combined.append(f"[TABLE] {table['description']}")
        
        # Add image descriptions
        for image in self.images:
            if image.get('description'):
                combined.append(f"[IMAGE] {image['description']}")
            if image.get('ocr_text'):
                combined.append(f"[OCR] {image['ocr_text']}")
        
        return "\n\n".join(filter(None, combined))
    
    def get_content_summary(self) -> Dict[str, Any]:
        """Get summary of extracted content"""
        return {
            "text_length": len(self.text),
            "image_count": len(self.images),
            "table_count": len(self.tables),
            "content_type": self.content_type,
            "extraction_time": self.extraction_time,
            "has_multimodal_content": len(self.images) > 0 or len(self.tables) > 0
        }

# ═══════════════════════════════════════════════════════════════════════════
# Abstract Base Extractor
# ═══════════════════════════════════════════════════════════════════════════

class BaseExtractor(ABC):
    """Abstract base class for file extractors"""
    
    @abstractmethod
    def can_extract(self, file_path: str) -> bool:
        """Check if this extractor can handle the file"""
        pass
    
    @abstractmethod
    def extract(self, file_path: str) -> ExtractedContent:
        """Extract content from the file"""
        pass
    
    def _get_file_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract common file metadata"""
        path = Path(file_path)
        stat = path.stat()
        
        return {
            "filename": path.name,
            "file_size": stat.st_size,
            "file_extension": path.suffix.lower(),
            "created_time": stat.st_ctime,
            "modified_time": stat.st_mtime,
            "file_type": self._detect_file_type(file_path)
        }
    
    def _detect_file_type(self, file_path: str) -> str:
        """Detect file type from extension"""
        ext = Path(file_path).suffix.lower()
        type_mapping = {
            '.pdf': 'document',
            '.docx': 'document', '.doc': 'document',
            '.xlsx': 'spreadsheet', '.xls': 'spreadsheet',
            '.pptx': 'presentation', '.ppt': 'presentation',
            '.jpg': 'image', '.jpeg': 'image', '.png': 'image', 
            '.tiff': 'image', '.bmp': 'image', '.gif': 'image',
            '.txt': 'text', '.rtf': 'text', '.md': 'text'
        }
        return type_mapping.get(ext, 'unknown')

# ═══════════════════════════════════════════════════════════════════════════
# PDF Extractor with Advanced Features
# ═══════════════════════════════════════════════════════════════════════════

class AdvancedPDFExtractor(BaseExtractor):
    """Advanced PDF extractor with image and table support"""
    
    def can_extract(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.pdf'
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        text_content = []
        images = []
        tables = []
        
        try:
            # Primary extraction with pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                    
                    # Extract tables
                    page_tables = page.extract_tables()
                    for table_idx, table in enumerate(page_tables):
                        if table:
                            table_df = pd.DataFrame(table[1:], columns=table[0])
                            table_description = self._describe_table(table_df)
                            tables.append({
                                'page': page_num + 1,
                                'table_index': table_idx,
                                'data': table_df.to_dict('records'),
                                'description': table_description,
                                'csv_representation': table_df.to_csv(index=False)
                            })
            
            # Advanced image extraction with PyMuPDF if available
            if PYMUPDF_AVAILABLE:
                images = self._extract_images_pymupdf(file_path)
            
        except Exception as e:
            print(f"PDF extraction error: {e}")
            # Fallback to simple text extraction
            try:
                with pdfplumber.open(file_path) as pdf:
                    text_content = [page.extract_text() or "" for page in pdf.pages]
            except Exception as fallback_error:
                print(f"PDF fallback extraction failed: {fallback_error}")
                text_content = ["Error extracting PDF content"]
        
        combined_text = "\n\n".join(filter(None, text_content))
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=combined_text,
            images=images,
            tables=tables,
            metadata=self._get_file_metadata(file_path),
            content_type="pdf_document",
            file_path=file_path,
            extraction_time=extraction_time
        )
    
    def _extract_images_pymupdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract images using PyMuPDF"""
        images = []
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        # Perform OCR if image contains text
                        ocr_text = self._perform_ocr(img_data)
                        
                        images.append({
                            'page': page_num + 1,
                            'image_index': img_index,
                            'format': 'png',
                            'size': len(img_data),
                            'ocr_text': ocr_text,
                            'description': f"Image from page {page_num + 1}"
                        })
                    
                    pix = None
            doc.close()
        except Exception as e:
            print(f"Image extraction error: {e}")
        
        return images
    
    def _perform_ocr(self, image_data: bytes) -> str:
        """Perform OCR on image data"""
        try:
            from io import BytesIO
            image = Image.open(BytesIO(image_data))
            ocr_text = pytesseract.image_to_string(image)
            return ocr_text.strip()
        except Exception as e:
            print(f"OCR error: {e}")
            return ""
    
    def _describe_table(self, df: pd.DataFrame) -> str:
        """Generate a description of the table"""
        try:
            rows, cols = df.shape
            columns = list(df.columns)
            description = f"Table with {rows} rows and {cols} columns. Columns: {', '.join(columns[:5])}"
            if len(columns) > 5:
                description += f" and {len(columns) - 5} more"
            return description
        except Exception:
            return "Table data"

# ═══════════════════════════════════════════════════════════════════════════
# Word Document Extractor
# ═══════════════════════════════════════════════════════════════════════════

class WordDocumentExtractor(BaseExtractor):
    """Word document extractor with table and image support"""
    
    def can_extract(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in {'.docx', '.doc'}
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        text_content = ""
        images = []
        tables = []
        
        try:
            if DOCX_ADVANCED_AVAILABLE and Path(file_path).suffix.lower() == '.docx':
                # Advanced extraction with python-docx
                doc = DocxDocument(file_path)
                
                # Extract paragraphs
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text_content = "\n\n".join(paragraphs)
                
                # Extract tables
                for table_idx, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        tables.append({
                            'table_index': table_idx,
                            'data': df.to_dict('records'),
                            'description': self._describe_table(df),
                            'csv_representation': df.to_csv(index=False)
                        })
                
                # Note: Image extraction from DOCX requires additional libraries
                # For now, we'll note their presence
                images = [{'description': 'Images present in document (extraction not implemented)'}]
                
            else:
                # Fallback to docx2txt
                text_content = docx2txt.process(file_path)
                
        except Exception as e:
            print(f"Word document extraction error: {e}")
            try:
                text_content = docx2txt.process(file_path)
            except Exception as fallback_error:
                print(f"Word document fallback failed: {fallback_error}")
                text_content = "Error extracting Word document content"
        
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=text_content,
            images=images,
            tables=tables,
            metadata=self._get_file_metadata(file_path),
            content_type="word_document",
            file_path=file_path,
            extraction_time=extraction_time
        )
    
    def _describe_table(self, df: pd.DataFrame) -> str:
        """Generate a description of the table"""
        try:
            rows, cols = df.shape
            columns = list(df.columns)
            description = f"Table with {rows} rows and {cols} columns. Columns: {', '.join(columns[:5])}"
            if len(columns) > 5:
                description += f" and {len(columns) - 5} more"
            return description
        except Exception:
            return "Table data"

# ═══════════════════════════════════════════════════════════════════════════
# Excel Spreadsheet Extractor
# ═══════════════════════════════════════════════════════════════════════════

class ExcelExtractor(BaseExtractor):
    """Excel spreadsheet extractor with sheet and table analysis"""
    
    def can_extract(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in {'.xlsx', '.xls'}
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        text_content = []
        tables = []
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Generate text description
                sheet_description = f"[Sheet: {sheet_name}]\n"
                sheet_description += f"Dimensions: {df.shape[0]} rows × {df.shape[1]} columns\n"
                sheet_description += f"Columns: {', '.join(df.columns.astype(str))}\n"
                
                # Add sample data
                if not df.empty:
                    sheet_description += "\nSample data:\n"
                    sheet_description += df.head(3).to_string(index=False)
                
                text_content.append(sheet_description)
                
                # Store table data
                tables.append({
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'description': f"Excel sheet '{sheet_name}' with {df.shape[0]} rows and {df.shape[1]} columns",
                    'csv_representation': df.to_csv(index=False),
                    'summary_stats': self._get_sheet_summary(df)
                })
                
        except Exception as e:
            print(f"Excel extraction error: {e}")
            try:
                # Fallback to simple CSV conversion
                dfs = pd.read_excel(file_path, sheet_name=None)
                text_content = [df.to_csv(index=False) for df in dfs.values()]
            except Exception as fallback_error:
                print(f"Excel fallback failed: {fallback_error}")
                text_content = ["Error extracting Excel content"]
        
        combined_text = "\n\n".join(text_content)
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=combined_text,
            images=[],  # Excel doesn't typically contain images we can extract
            tables=tables,
            metadata=self._get_file_metadata(file_path),
            content_type="excel_spreadsheet",
            file_path=file_path,
            extraction_time=extraction_time
        )
    
    def _get_sheet_summary(self, df: pd.DataFrame) -> Dict[str, Any]:
        """Generate summary statistics for a sheet"""
        try:
            numeric_cols = df.select_dtypes(include=['number']).columns
            summary = {
                'total_rows': len(df),
                'total_columns': len(df.columns),
                'numeric_columns': len(numeric_cols),
                'text_columns': len(df.columns) - len(numeric_cols),
                'empty_cells': df.isnull().sum().sum()
            }
            
            if len(numeric_cols) > 0:
                summary['numeric_summary'] = df[numeric_cols].describe().to_dict()
            
            return summary
        except Exception:
            return {'error': 'Could not generate summary'}

# ═══════════════════════════════════════════════════════════════════════════
# PowerPoint Extractor
# ═══════════════════════════════════════════════════════════════════════════

class PowerPointExtractor(BaseExtractor):
    """PowerPoint presentation extractor"""
    
    def can_extract(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in {'.pptx', '.ppt'} and PPTX_AVAILABLE
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        text_content = []
        images = []
        tables = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = f"[Slide {slide_idx + 1}]\n"
                slide_content = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # Check for tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data:
                            df = pd.DataFrame(table_data[1:], columns=table_data[0])
                            tables.append({
                                'slide': slide_idx + 1,
                                'data': df.to_dict('records'),
                                'description': f"Table from slide {slide_idx + 1}",
                                'csv_representation': df.to_csv(index=False)
                            })
                
                if slide_content:
                    slide_text += "\n".join(slide_content)
                    text_content.append(slide_text)
                
                # Note: Image extraction from PPTX requires additional processing
                # For now, we'll count them
                image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # Picture type
                if image_count > 0:
                    images.append({
                        'slide': slide_idx + 1,
                        'count': image_count,
                        'description': f"{image_count} images on slide {slide_idx + 1}"
                    })
                    
        except Exception as e:
            print(f"PowerPoint extraction error: {e}")
            text_content = ["Error extracting PowerPoint content"]
        
        combined_text = "\n\n".join(text_content)
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=combined_text,
            images=images,
            tables=tables,
            metadata=self._get_file_metadata(file_path),
            content_type="powerpoint_presentation",
            file_path=file_path,
            extraction_time=extraction_time
        )

# ═══════════════════════════════════════════════════════════════════════════
# Image Extractor with OCR
# ═══════════════════════════════════════════════════════════════════════════

class ImageExtractor(BaseExtractor):
    """Image extractor with OCR capabilities"""
    
    def can_extract(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in {'.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.gif'}
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        text_content = ""
        images = []
        
        try:
            # Perform OCR on the image
            image = Image.open(file_path)
            ocr_text = pytesseract.image_to_string(image)
            
            text_content = f"[OCR Text from {Path(file_path).name}]\n{ocr_text.strip()}"
            
            # Store image metadata
            images.append({
                'filename': Path(file_path).name,
                'format': image.format,
                'size': image.size,
                'mode': image.mode,
                'ocr_text': ocr_text.strip(),
                'description': f"Image file with OCR text extraction"
            })
            
        except Exception as e:
            print(f"Image extraction error: {e}")
            text_content = f"Error extracting text from image: {Path(file_path).name}"
        
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=text_content,
            images=images,
            tables=[],
            metadata=self._get_file_metadata(file_path),
            content_type="image_file",
            file_path=file_path,
            extraction_time=extraction_time
        )

# ═══════════════════════════════════════════════════════════════════════════
# Text File Extractor
# ═══════════════════════════════════════════════════════════════════════════

class TextFileExtractor(BaseExtractor):
    """Simple text file extractor"""
    
    def can_extract(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in {'.txt', '.rtf', '.md', '.csv'}
    
    def extract(self, file_path: str) -> ExtractedContent:
        start_time = time.time()
        
        try:
            text_content = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        except Exception as e:
            print(f"Text file extraction error: {e}")
            text_content = "Error reading text file"
        
        extraction_time = time.time() - start_time
        
        return ExtractedContent(
            text=text_content,
            images=[],
            tables=[],
            metadata=self._get_file_metadata(file_path),
            content_type="text_file",
            file_path=file_path,
            extraction_time=extraction_time
        )

# ═══════════════════════════════════════════════════════════════════════════
# Multimodal Extraction Manager
# ═══════════════════════════════════════════════════════════════════════════

class MultimodalExtractionManager:
    """Central manager for multimodal file extraction"""
    
    def __init__(self):
        self.extractors = [
            AdvancedPDFExtractor(),
            WordDocumentExtractor(),
            ExcelExtractor(),
            PowerPointExtractor() if PPTX_AVAILABLE else None,
            ImageExtractor(),
            TextFileExtractor()
        ]
        # Filter out None extractors
        self.extractors = [e for e in self.extractors if e is not None]
        
        self.extraction_stats = {
            'total_extractions': 0,
            'successful_extractions': 0,
            'failed_extractions': 0,
            'total_extraction_time': 0,
            'by_file_type': {}
        }
    
    def extract_content(self, file_path: str) -> ExtractedContent:
        """Extract content using the appropriate extractor"""
        start_time = time.time()
        self.extraction_stats['total_extractions'] += 1
        
        # Find appropriate extractor
        extractor = None
        for ext in self.extractors:
            if ext.can_extract(file_path):
                extractor = ext
                break
        
        if not extractor:
            # Fallback to text extraction
            print(f"No specific extractor found for {file_path}, using text fallback")
            extractor = TextFileExtractor()
        
        try:
            content = extractor.extract(file_path)
            self.extraction_stats['successful_extractions'] += 1
            
            # Update stats
            file_type = content.content_type
            if file_type not in self.extraction_stats['by_file_type']:
                self.extraction_stats['by_file_type'][file_type] = 0
            self.extraction_stats['by_file_type'][file_type] += 1
            
            total_time = time.time() - start_time
            self.extraction_stats['total_extraction_time'] += total_time
            
            print(f"✅ Extracted {file_type}: {content.get_content_summary()}")
            return content
            
        except Exception as e:
            self.extraction_stats['failed_extractions'] += 1
            print(f"❌ Extraction failed for {file_path}: {e}")
            
            # Return error content
            return ExtractedContent(
                text=f"Error extracting content from {Path(file_path).name}: {str(e)}",
                images=[],
                tables=[],
                metadata={'error': str(e)},
                content_type="error",
                file_path=file_path,
                extraction_time=time.time() - start_time
            )
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        formats = []
        for extractor in self.extractors:
            if isinstance(extractor, AdvancedPDFExtractor):
                formats.extend(['.pdf'])
            elif isinstance(extractor, WordDocumentExtractor):
                formats.extend(['.docx', '.doc'])
            elif isinstance(extractor, ExcelExtractor):
                formats.extend(['.xlsx', '.xls'])
            elif isinstance(extractor, PowerPointExtractor):
                formats.extend(['.pptx', '.ppt'])
            elif isinstance(extractor, ImageExtractor):
                formats.extend(['.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.gif'])
            elif isinstance(extractor, TextFileExtractor):
                formats.extend(['.txt', '.rtf', '.md', '.csv'])
        
        return sorted(list(set(formats)))
    
    def get_extraction_stats(self) -> Dict[str, Any]:
        """Get extraction statistics"""
        total = self.extraction_stats['total_extractions']
        avg_time = (self.extraction_stats['total_extraction_time'] / total 
                   if total > 0 else 0)
        
        return {
            **self.extraction_stats,
            'success_rate': (self.extraction_stats['successful_extractions'] / total 
                           if total > 0 else 0),
            'average_extraction_time': avg_time,
            'supported_formats': self.get_supported_formats()
        }

# ═══════════════════════════════════════════════════════════════════════════
# Global Instance and Backward Compatibility
# ═══════════════════════════════════════════════════════════════════════════

# Global multimodal extraction manager
multimodal_manager = MultimodalExtractionManager()

def extract_text_multimodal(file_path: str) -> str:
    """
    Backward compatible function that returns combined text
    This maintains compatibility with existing pipeline
    """
    content = multimodal_manager.extract_content(file_path)
    return content.get_combined_text()

def extract_content_advanced(file_path: str) -> ExtractedContent:
    """
    Advanced extraction function that returns full multimodal content
    Use this for new implementations that need multimodal data
    """
    return multimodal_manager.extract_content(file_path)

# Backward compatibility alias
extract_text = extract_text_multimodal
